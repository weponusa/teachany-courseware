#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
TeachAny · HTML → PPTX 导出工具（v5.34.12 · 图文美观版）

用法:
    python3 scripts/export-pptx.py <课件目录>

行为（v5.34.12 升级）:
    - 真正的"图文 PPTX"：封面大图 hero / 节首图卡 / 双栏结构化 bullet /
      卡片式练习题 / 互动占位大图标
    - 配色：浅色主题（白底 + 品牌紫 + 琥珀点缀），不再是"顶部大蓝块 + 左文右图"
      的办公风
    - 不调用任何第三方生图 API。图片复用课件 assets/ 下 image_gen（WorkBuddy
      原生工具）已生成的，或 HTML 中 <img> 本地引用的
    - 若 section 无图 → 自动选择"纯文字特制版式"（大字号标题 + 彩色装饰 +
      引号大卡片），而不是留空白
    - 产出: <课件目录>/<课件名>.pptx

依赖:
    pip3 install python-pptx beautifulsoup4

⚠️ 图片来源承诺（v5.34.12）:
    本脚本不持有任何第三方生图服务的密钥，不会调用 Gemini / OpenAI DALL-E /
    Replicate 等接口。所有图片必须由 AI 在 HTML 生成阶段通过 `image_gen`
    （WorkBuddy / CodeBuddy 等宿主 IDE 原生提供的工具）预先生成并存入
    课件 assets/ 目录。PPTX 是 HTML 的派生件，不单独生图。

遵循硬规则 #46 / #47 PPTX 导出基线 (TeachAny SKILL_CN.md)。
"""

import argparse
import json
import os
import re
import sys
from pathlib import Path

# ───────────────────────────── 依赖自动安装 ─────────────────────────────

def _ensure_deps():
    missing = []
    try:
        import pptx  # noqa: F401
    except ImportError:
        missing.append('python-pptx')
    try:
        from bs4 import BeautifulSoup  # noqa: F401
    except ImportError:
        missing.append('beautifulsoup4')

    if missing:
        print(f"🔧 检测到缺失依赖: {', '.join(missing)}，尝试自动安装...")
        import subprocess
        ret = subprocess.call(
            [sys.executable, '-m', 'pip', 'install', '--quiet', *missing]
        )
        if ret != 0:
            print(f"❌ 自动安装失败。请手动执行: pip3 install {' '.join(missing)}")
            sys.exit(2)
        print("✅ 依赖安装完成")


_ensure_deps()

from pptx import Presentation  # noqa: E402
from pptx.util import Inches, Pt, Emu  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR  # noqa: E402
from pptx.oxml.ns import qn  # noqa: E402
from lxml import etree  # noqa: E402
from bs4 import BeautifulSoup  # noqa: E402


# ───────────────────────────── 设计 Token ─────────────────────────────

SLIDE_W = Inches(13.333)   # 16:9
SLIDE_H = Inches(7.5)

# 浅色主题（白底 + 品牌紫 + 琥珀点缀 + 深灰）
THEME = {
    'bg':         RGBColor(0xff, 0xff, 0xff),
    'bg_soft':    RGBColor(0xf8, 0xfa, 0xfc),
    'bg_card':    RGBColor(0xfd, 0xfd, 0xff),
    'primary':    RGBColor(0x63, 0x66, 0xf1),    # indigo-500
    'primary_dk': RGBColor(0x4f, 0x46, 0xe5),    # indigo-600
    'primary_lt': RGBColor(0xe0, 0xe7, 0xff),    # indigo-100
    'accent':     RGBColor(0xf5, 0x9e, 0x0b),    # amber-500
    'accent_lt':  RGBColor(0xfe, 0xf3, 0xc7),    # amber-100
    'success':    RGBColor(0x10, 0xb9, 0x81),    # emerald-500
    'text_1':     RGBColor(0x0f, 0x17, 0x2a),    # slate-900
    'text_2':     RGBColor(0x33, 0x41, 0x55),    # slate-700
    'text_3':     RGBColor(0x64, 0x74, 0x8b),    # slate-500
    'border':     RGBColor(0xe2, 0xe8, 0xf0),    # slate-200
    'border_soft': RGBColor(0xf1, 0xf5, 0xf9),    # slate-100
}

# 字体（中文 + 英文双栈）
FONT_CN = 'Microsoft YaHei'
FONT_EN = 'Inter'

# 互动组件降级图标（只保留最关键的 5 种）
INTERACTIVE_MARKERS = [
    ('<canvas',        '🎨', 'Canvas 互动'),
    ('knowledge-graph', '🗺️', '知识图谱'),
    ('audioPlaylist',  '🔊', '语音讲解'),
    ('ai-tutor',       '🤖', 'AI 学伴'),
    ('<video',         '🎬', '教学视频'),
]


# ───────────────────────────── 工具函数 ─────────────────────────────

def clean_text(s: str) -> str:
    if not s:
        return ''
    s = re.sub(r'\s+', ' ', s)
    s = s.replace('\u200b', '').replace('\ufeff', '').strip()
    return s


def read_manifest(course_dir: Path) -> dict:
    p = course_dir / 'manifest.json'
    if not p.exists():
        return {}
    try:
        return json.loads(p.read_text(encoding='utf-8'))
    except Exception as e:
        print(f"⚠️  manifest 解析失败: {e}")
        return {}


def find_images_in_section(section, course_dir: Path, max_n=2):
    """收集 section 里的本地图（按出现顺序），最多 max_n 张"""
    out = []
    for img in section.find_all('img'):
        if len(out) >= max_n:
            break
        src = (img.get('src') or '').strip()
        if not src or src.startswith(('data:', 'http://', 'https://')):
            continue
        candidate = (course_dir / src).resolve()
        try:
            candidate.relative_to(course_dir.resolve())
        except ValueError:
            continue
        if candidate.exists() and candidate.suffix.lower() in (
            '.png', '.jpg', '.jpeg', '.gif', '.webp'
        ):
            out.append(candidate)
    return out


def course_hero_image(course_dir: Path):
    """封面图：统一检测逻辑（v6.2 规范）
    优先级：
      1. *-hero.{png,jpg,webp}   后缀匹配（主流模式）
      2. hero-*.{png,jpg,webp}   前缀匹配（兼容旧命名）
      3. hero.{png,jpg,webp}     纯名称匹配
      4. 图片目录下第一张图       兜底
    搜索 assets/ 和 images/ 两个目录
    """
    img_exts = ('.png', '.jpg', '.jpeg', '.webp')
    img_dir = None
    for name in ('assets', 'images'):
        candidate = course_dir / name
        if candidate.exists() and candidate.is_dir():
            img_dir = candidate
            break
    if img_dir is None:
        return None
    all_imgs = sorted([
        p for p in img_dir.iterdir()
        if p.is_file() and p.suffix.lower() in img_exts
    ])
    if not all_imgs:
        return None
    # 1. 后缀匹配：*-hero.ext（主流模式，如 linear-hero.png）
    for p in all_imgs:
        if p.stem.lower().endswith('-hero'):
            return p
    # 2. 前缀匹配：hero-*.ext（兼容旧命名，如 hero-denglouque.png）
    for p in all_imgs:
        if p.stem.lower().startswith('hero-'):
            return p
    # 3. 纯名称匹配：hero.ext
    for p in all_imgs:
        if p.stem.lower() == 'hero':
            return p
    # 4. 兜底：第一张图
    return all_imgs[0] if all_imgs else None


def detect_interactive_components(section):
    html_str = str(section).lower()
    found = []
    for marker, icon, label in INTERACTIVE_MARKERS:
        if marker.lower() in html_str:
            found.append((icon, label))
    return found


def extract_section_texts(section):
    """从 section 抽文本，区分 bullet / 段落 / 题目"""
    title = ''
    for tag in ['h1', 'h2', 'h3']:
        h = section.find(tag)
        if h:
            title = clean_text(h.get_text())
            break

    # bullet 优先从 li 抽，再 fallback 到 p
    bullets = []
    for li in section.find_all('li'):
        t = clean_text(li.get_text())
        if t and t != title and len(t) > 2:
            bullets.append(t)

    # 段落（bullet 不足时用 p 补）
    paragraphs = []
    for p in section.find_all('p'):
        t = clean_text(p.get_text())
        if t and t != title and len(t) > 2:
            paragraphs.append(t)

    # 拼：bullet 太少时把 paragraph 当 bullet 用
    if len(bullets) < 3 and paragraphs:
        bullets = paragraphs + bullets  # paragraph 在前，通常更 important

    # 题目识别（保留做卡片）
    quizzes = []
    for q in section.find_all(attrs={'class': re.compile(r'quiz|exercise|problem')}):
        t = clean_text(q.get_text())
        if t and len(t) > 5:
            quizzes.append(t[:240])

    # 也从 data-question 抽
    for q in section.find_all(attrs={'data-question': True}):
        t = clean_text(q['data-question'])
        if t:
            quizzes.append(t[:240])

    return {
        'title':    title,
        'bullets':  bullets[:8],          # 最多 8 条
        'quizzes':  quizzes[:4],          # 最多 4 题
        'kicker':   _derive_kicker(section),
    }


def _derive_kicker(section):
    """从 section id / data-module 中推导小标签（如 '模块一' / '前测' / '综合任务'）"""
    sid = (section.get('id') or '').lower()
    kmap = [
        ('hero',        '封面'),
        ('objective',   '学习目标'),
        ('pretest',     '前测诊断'),
        ('post',        '后测挑战'),
        ('summary',     '课堂小结'),
        ('module1',     '模块一'), ('m1', '模块一'),
        ('module2',     '模块二'), ('m2', '模块二'),
        ('module3',     '模块三'), ('m3', '模块三'),
        ('module4',     '模块四'), ('m4', '模块四'),
        ('quiz',        '练习'),
        ('comprehensive', '综合任务'),
        ('knowledge-graph', '知识图谱'),
        ('audio',       '语音讲解'),
    ]
    for k, v in kmap:
        if k in sid:
            return v
    return ''


# ───────────────────────────── 底层绘图原语 ─────────────────────────────

def add_rect(slide, left, top, w, h, fill_rgb, line_rgb=None, line_w=None):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill_rgb
    if line_rgb:
        rect.line.color.rgb = line_rgb
        if line_w:
            rect.line.width = line_w
    else:
        rect.line.fill.background()
    rect.shadow.inherit = False
    return rect


def add_rounded_rect(slide, left, top, w, h, fill_rgb, line_rgb=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    # 调整圆角比例
    try:
        shape.adjustments[0] = 0.08
    except Exception:
        pass
    return shape


def add_text(slide, left, top, w, h, text, *,
             font_size=14, bold=False, color=None, font=FONT_CN,
             align='left', anchor='top', line_spacing=1.2):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    anchors = {'top': MSO_ANCHOR.TOP, 'middle': MSO_ANCHOR.MIDDLE, 'bottom': MSO_ANCHOR.BOTTOM}
    tf.vertical_anchor = anchors.get(anchor, MSO_ANCHOR.TOP)
    aligns = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}
    p = tf.paragraphs[0]
    p.alignment = aligns.get(align, PP_ALIGN.LEFT)
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.name = font
    # 双字体回退（英文用 Inter，中文用雅黑）
    rPr = r._r.get_or_add_rPr()
    ea = etree.SubElement(rPr, qn('a:ea'))
    ea.set('typeface', FONT_CN)
    latin = etree.SubElement(rPr, qn('a:latin'))
    latin.set('typeface', FONT_EN)
    if color is not None:
        r.font.color.rgb = color
    return tb


def add_multi_line_text(slide, left, top, w, h, lines, *,
                        font_size=14, color=None, bullet_char='●',
                        line_spacing=1.35, bullet_color=None, bold_first=False):
    """带 bullet 的多行文本"""
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(4)
        # bullet 字符
        if bullet_char:
            r_b = p.add_run()
            r_b.text = f'{bullet_char}  '
            r_b.font.size = Pt(font_size)
            r_b.font.color.rgb = bullet_color or THEME['primary']
            r_b.font.name = FONT_CN
        r = p.add_run()
        r.text = line
        r.font.size = Pt(font_size)
        r.font.name = FONT_CN
        r.font.bold = bold_first and i == 0
        if color is not None:
            r.font.color.rgb = color
    return tb


def add_image_safely(slide, path, left, top, w, h):
    """嵌图并按比例裁剪/缩放，超出时等比内切"""
    try:
        slide.shapes.add_picture(str(path), left, top, width=w, height=h)
        return True
    except Exception as e:
        print(f"⚠️  图片嵌入失败 {path.name}: {e}")
        return False


def new_blank_slide(pres):
    # layouts[6] 通常是 blank
    return pres.slides.add_slide(pres.slide_layouts[6])


# ───────────────────────────── 版式 ─────────────────────────────

def add_title_slide(pres, course_title, subtitle, hero_image: Path = None):
    """封面：左 45% 文字 + 右 55% hero 图（无图时整栏品牌色装饰）"""
    slide = new_blank_slide(pres)
    # 整体白底
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, THEME['bg'])

    # 右侧 hero 图区
    right_x = Inches(6.2)
    right_w = SLIDE_W - right_x
    if hero_image and hero_image.exists():
        # 图填满右侧，外层加一个浅色底
        add_rect(slide, right_x, 0, right_w, SLIDE_H, THEME['primary_lt'])
        add_image_safely(slide, hero_image, right_x, 0, right_w, SLIDE_H)
    else:
        # 无图时用品牌色色块 + 装饰圆点
        add_rect(slide, right_x, 0, right_w, SLIDE_H, THEME['primary'])
        for i, (cx, cy, r_in) in enumerate([
            (right_x + Inches(2), Inches(1.5), Inches(1.0)),
            (right_x + Inches(4), Inches(4), Inches(1.6)),
            (right_x + Inches(1), Inches(5.5), Inches(0.8)),
        ]):
            circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, r_in, r_in)
            circle.fill.solid()
            circle.fill.fore_color.rgb = THEME['accent'] if i == 1 else THEME['primary_dk']
            circle.line.fill.background()

    # 左侧文字区
    pad_x = Inches(0.8)
    # 顶部小 Kicker
    add_text(slide, pad_x, Inches(1.0), Inches(5), Inches(0.4),
             'TeachAny · 教学课件',
             font_size=14, color=THEME['accent'], bold=True)
    # 彩色装饰条
    add_rect(slide, pad_x, Inches(1.55), Inches(0.6), Inches(0.08), THEME['accent'])

    # 大标题
    add_text(slide, pad_x, Inches(1.9), Inches(5.3), Inches(2.6),
             course_title,
             font_size=40, bold=True, color=THEME['text_1'], line_spacing=1.15)
    # 副标题
    if subtitle:
        add_text(slide, pad_x, Inches(4.7), Inches(5.3), Inches(0.6),
                 subtitle,
                 font_size=16, color=THEME['text_3'])
    # 底部徽章
    add_rect(slide, pad_x, SLIDE_H - Inches(1.1), Inches(4.5), Inches(0.4), THEME['primary_lt'])
    add_text(slide, pad_x + Inches(0.15), SLIDE_H - Inches(1.05), Inches(4.2), Inches(0.35),
             '⚡ 互动 HTML 课件  ·  PPTX 派生版',
             font_size=12, color=THEME['primary_dk'], bold=True, anchor='middle')


def add_section_hero_slide(pres, kicker: str, title: str, image: Path = None):
    """节首图卡：模块转场时用，大字号标题 + 配图"""
    slide = new_blank_slide(pres)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, THEME['bg'])

    # 左彩色装饰条（2/3 slide 高）
    add_rect(slide, 0, Inches(1.2), Inches(0.15), Inches(5.1), THEME['accent'])

    # 左文字
    text_x = Inches(0.7)
    text_w = Inches(7.0)
    if kicker:
        add_text(slide, text_x, Inches(1.3), text_w, Inches(0.5),
                 kicker, font_size=15, color=THEME['accent'], bold=True)
    add_text(slide, text_x, Inches(2.0), text_w, Inches(3.0),
             title or '', font_size=44, bold=True,
             color=THEME['text_1'], line_spacing=1.1)
    # 底部 progress bar（装饰）
    add_rect(slide, text_x, Inches(6.0), Inches(4), Inches(0.06), THEME['primary_lt'])
    add_rect(slide, text_x, Inches(6.0), Inches(1.2), Inches(0.06), THEME['primary'])

    # 右图（卡片式圆角 + 背景）
    if image and image.exists():
        img_x = Inches(8.2)
        img_w = Inches(4.5)
        img_y = Inches(1.3)
        img_h = Inches(5.0)
        add_rounded_rect(slide, img_x - Inches(0.1), img_y - Inches(0.1),
                         img_w + Inches(0.2), img_h + Inches(0.2),
                         THEME['primary_lt'])
        add_image_safely(slide, image, img_x, img_y, img_w, img_h)


def add_content_slide(pres, title: str, bullets: list, image: Path = None,
                      kicker: str = '', interactive: list = None, idx: int = 0):
    """标准内容页：标题 + 双栏 bullet + 可选配图 + kicker + 页码装饰"""
    slide = new_blank_slide(pres)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, THEME['bg'])

    # 顶部细装饰条
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), THEME['primary'])

    # Kicker + 大标题
    top_x = Inches(0.7)
    if kicker:
        add_text(slide, top_x, Inches(0.35), Inches(4), Inches(0.4),
                 kicker, font_size=13, color=THEME['accent'], bold=True)
    add_text(slide, top_x, Inches(0.75), Inches(12), Inches(0.9),
             title or '(无标题)',
             font_size=28, bold=True, color=THEME['text_1'])
    # 标题下细短线
    add_rect(slide, top_x, Inches(1.6), Inches(0.5), Inches(0.05), THEME['accent'])

    # 内容区布局
    has_image = image and image.exists()
    content_top = Inches(2.0)
    content_h = Inches(5.0)

    # 按 bullet 数量决定排版策略
    n = len(bullets)
    if has_image:
        # 图在右，文本在左
        text_x = top_x
        text_w = Inches(7.2)
        _render_bullet_card(slide, text_x, content_top, text_w, content_h, bullets)

        img_x = Inches(8.3)
        img_w = Inches(4.3)
        img_y = content_top
        img_h = content_h
        add_rounded_rect(slide, img_x - Inches(0.08), img_y - Inches(0.08),
                         img_w + Inches(0.16), img_h + Inches(0.16),
                         THEME['primary_lt'])
        add_image_safely(slide, image, img_x, img_y, img_w, img_h)
    elif n >= 4:
        # 无图：双栏 bullet
        col_w = Inches(5.8)
        gap = Inches(0.4)
        mid = n // 2 + (n % 2)
        _render_bullet_card(slide, top_x, content_top, col_w, content_h, bullets[:mid])
        _render_bullet_card(slide, top_x + col_w + gap, content_top, col_w, content_h, bullets[mid:])
    else:
        # 无图 + bullet 少：大字号放大展示
        _render_highlight_card(slide, top_x, content_top, SLIDE_W - Inches(1.4), content_h, bullets)

    # 右下角装饰（页码 + slide 号）
    add_text(slide, SLIDE_W - Inches(1.5), SLIDE_H - Inches(0.5),
             Inches(1.2), Inches(0.3),
             f'— {idx:02d} —', font_size=11, color=THEME['text_3'],
             align='right')

    # 互动降级提示
    if interactive:
        note = '  '.join(f'{ic} {lb}' for ic, lb in interactive[:4])
        add_text(slide, top_x, SLIDE_H - Inches(0.5), Inches(10), Inches(0.3),
                 f'🔗 含互动：{note}（请在 HTML 课件中体验）',
                 font_size=10, color=THEME['text_3'])


def _render_bullet_card(slide, x, y, w, h, bullets):
    """bullet 卡片：每条带数字徽章"""
    bullets = [b[:120] for b in bullets if b]
    if not bullets:
        return
    # 背景轻色卡片
    add_rounded_rect(slide, x, y, w, h, THEME['bg_soft'],
                     line_rgb=THEME['border_soft'])
    # 内容
    inner_x = x + Inches(0.3)
    inner_y = y + Inches(0.25)
    line_h = Inches(0.7)
    for i, text in enumerate(bullets[:6]):
        # 数字徽章
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                       inner_x, inner_y + line_h * i + Inches(0.06),
                                       Inches(0.35), Inches(0.35))
        badge.fill.solid()
        badge.fill.fore_color.rgb = THEME['primary']
        badge.line.fill.background()
        badge.shadow.inherit = False
        add_text(slide, inner_x, inner_y + line_h * i + Inches(0.06),
                 Inches(0.35), Inches(0.35),
                 str(i + 1), font_size=12, bold=True,
                 color=THEME['bg'], align='center', anchor='middle')
        # 文本
        add_text(slide, inner_x + Inches(0.5),
                 inner_y + line_h * i,
                 w - Inches(0.9), line_h,
                 text, font_size=15, color=THEME['text_2'],
                 anchor='middle', line_spacing=1.25)


def _render_highlight_card(slide, x, y, w, h, bullets):
    """高亮卡片：bullet 很少时用大字号 + 引号装饰"""
    if not bullets:
        return
    # 大引号装饰
    add_text(slide, x + Inches(0.3), y + Inches(0.1), Inches(1.0), Inches(1.0),
             '"', font_size=96, color=THEME['primary_lt'], bold=True, font=FONT_EN)
    inner_y = y + Inches(0.8)
    line_h = Inches(1.0)
    for i, text in enumerate(bullets[:3]):
        add_text(slide, x + Inches(0.9), inner_y + line_h * i,
                 w - Inches(1.0), line_h,
                 text, font_size=22, bold=(i == 0),
                 color=THEME['text_1'] if i == 0 else THEME['text_2'],
                 anchor='middle', line_spacing=1.3)


def add_quiz_cards_slide(pres, section_title: str, quizzes: list, kicker: str = ''):
    """练习题卡片页：最多 2-3 题以卡片形式展示"""
    slide = new_blank_slide(pres)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, THEME['bg'])
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), THEME['accent'])

    top_x = Inches(0.7)
    if kicker:
        add_text(slide, top_x, Inches(0.35), Inches(4), Inches(0.4),
                 kicker, font_size=13, color=THEME['accent'], bold=True)
    add_text(slide, top_x, Inches(0.75), Inches(12), Inches(0.9),
             f'📝 {section_title or "课堂练习"}',
             font_size=26, bold=True, color=THEME['text_1'])

    # 卡片布局
    card_y = Inches(2.0)
    card_h = Inches(4.8)
    count = min(len(quizzes), 3)
    card_w = (SLIDE_W - Inches(1.4) - Inches(0.4) * (count - 1)) / max(count, 1)
    for i, q in enumerate(quizzes[:count]):
        cx = top_x + (card_w + Inches(0.4)) * i
        add_rounded_rect(slide, cx, card_y, card_w, card_h, THEME['bg_card'],
                         line_rgb=THEME['border'])
        # 卡片题号徽章
        add_rect(slide, cx + Inches(0.3), card_y + Inches(0.3),
                 Inches(0.8), Inches(0.35), THEME['accent_lt'])
        add_text(slide, cx + Inches(0.3), card_y + Inches(0.3),
                 Inches(0.8), Inches(0.35),
                 f' Q{i + 1} ', font_size=12, bold=True,
                 color=THEME['accent'], align='center', anchor='middle')
        # 题干
        add_text(slide, cx + Inches(0.3), card_y + Inches(0.9),
                 card_w - Inches(0.6), card_h - Inches(1.3),
                 q[:240], font_size=14, color=THEME['text_2'], line_spacing=1.4)


def add_interactive_placeholder_slide(pres, section_title: str, components: list,
                                      html_link: str = ''):
    """互动 section 占位：全屏大图标 + 入口"""
    slide = new_blank_slide(pres)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, THEME['bg_soft'])

    # 中央卡片
    card_w = Inches(9)
    card_h = Inches(5.5)
    card_x = (SLIDE_W - card_w) / 2
    card_y = (SLIDE_H - card_h) / 2
    add_rounded_rect(slide, card_x, card_y, card_w, card_h, THEME['bg'],
                     line_rgb=THEME['border'])

    # 大图标
    icons = '  '.join(ic for ic, _ in components) or '🎯'
    add_text(slide, card_x, card_y + Inches(0.6), card_w, Inches(1.5),
             icons, font_size=72, align='center', anchor='middle')

    # 标题
    add_text(slide, card_x, card_y + Inches(2.3), card_w, Inches(0.7),
             section_title or '线上互动环节',
             font_size=28, bold=True, color=THEME['text_1'],
             align='center', anchor='middle')

    # 组件名称
    labels = '  ·  '.join(lb for _, lb in components) or '互动组件'
    add_text(slide, card_x, card_y + Inches(3.1), card_w, Inches(0.5),
             labels, font_size=14, color=THEME['text_3'],
             align='center', anchor='middle')

    # 操作提示
    add_rect(slide, card_x + Inches(1), card_y + Inches(4.0),
             card_w - Inches(2), Inches(0.5), THEME['primary_lt'])
    add_text(slide, card_x + Inches(1), card_y + Inches(4.0),
             card_w - Inches(2), Inches(0.5),
             '➜  在 HTML 课件中体验完整互动',
             font_size=14, bold=True, color=THEME['primary_dk'],
             align='center', anchor='middle')


def add_end_slide(pres, course_title: str):
    """结尾页：简洁的品牌收束"""
    slide = new_blank_slide(pres)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, THEME['primary'])
    # 装饰圆
    for cx, cy, r_in, rgb in [
        (Inches(1.5), Inches(1), Inches(1.5), THEME['primary_dk']),
        (Inches(10), Inches(5), Inches(2), THEME['primary_dk']),
        (Inches(11.5), Inches(1.5), Inches(0.8), THEME['accent']),
    ]:
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, r_in, r_in)
        circle.fill.solid()
        circle.fill.fore_color.rgb = rgb
        circle.line.fill.background()

    # Thank you
    add_text(slide, Inches(1), Inches(2.8), Inches(11), Inches(1.5),
             '谢谢观看',
             font_size=72, bold=True, color=THEME['bg'],
             align='center', anchor='middle')
    add_text(slide, Inches(1), Inches(4.3), Inches(11), Inches(0.6),
             course_title,
             font_size=20, color=THEME['primary_lt'],
             align='center', anchor='middle')
    add_text(slide, Inches(1), Inches(5.3), Inches(11), Inches(0.5),
             '⚡ TeachAny · 开源教学课件平台',
             font_size=14, bold=True, color=THEME['accent'],
             align='center', anchor='middle')


# ───────────────────────────── 主流程 ─────────────────────────────

def export_pptx(course_dir: Path, output: Path = None) -> Path:
    course_dir = course_dir.resolve()
    html_path = course_dir / 'index.html'
    if not html_path.exists():
        raise FileNotFoundError(f'课件目录缺少 index.html: {html_path}')

    html = html_path.read_text(encoding='utf-8')
    soup = BeautifulSoup(html, 'html.parser')
    manifest = read_manifest(course_dir)

    course_title = (manifest.get('name') or manifest.get('title') or
                    clean_text(soup.title.string if soup.title else '') or
                    course_dir.name)
    # 去掉 title 里的 TeachAny 版本后缀（封面另有徽章）
    course_title = re.sub(r'\s*·?\s*TeachAny\s*v?[\d.]+\s*$', '', course_title).strip()

    subtitle_meta = []
    if manifest.get('subject'):
        subtitle_meta.append(manifest['subject'])
    if manifest.get('grade'):
        subtitle_meta.append(f"G{manifest['grade']}")
    if manifest.get('teachany_version'):
        subtitle_meta.append(f"TeachAny v{manifest['teachany_version']}")
    subtitle = '  ·  '.join(subtitle_meta)

    pres = Presentation()
    pres.slide_width = SLIDE_W
    pres.slide_height = SLIDE_H

    # 封面
    hero = course_hero_image(course_dir)
    add_title_slide(pres, course_title, subtitle, hero_image=hero)

    # 遍历 section
    sections = soup.find_all('section')
    if not sections:
        sections = soup.select('div.slide')

    produced_image_ids = set()
    if hero:
        produced_image_ids.add(str(hero))

    slide_idx = 1
    prev_kicker = None
    for sec in sections:
        texts = extract_section_texts(sec)
        imgs = find_images_in_section(sec, course_dir, max_n=2)
        interactive = detect_interactive_components(sec)
        kicker = texts['kicker']

        # 完全空 section 跳过
        if not texts['title'] and not texts['bullets'] and not texts['quizzes'] and not imgs:
            continue

        # 如果有很强的互动组件而无实际内容 → 占位页
        if interactive and not texts['bullets'] and not texts['quizzes']:
            add_interactive_placeholder_slide(
                pres,
                texts['title'] or kicker or '线上互动环节',
                interactive
            )
            slide_idx += 1
            continue

        # 节首图卡（仅当 kicker 变化到新"模块"时插入，避免每 section 都插）
        picking_image = None
        for p in imgs:
            if str(p) not in produced_image_ids:
                picking_image = p
                break
        is_new_module = (kicker != prev_kicker and kicker
                         and kicker in ('模块一', '模块二', '模块三', '模块四'))
        if is_new_module and picking_image is not None:
            add_section_hero_slide(pres, kicker, texts['title'], picking_image)
            produced_image_ids.add(str(picking_image))
            # 再拿下一张图给下面的内容页
            picking_image = None
            for p in imgs:
                if str(p) not in produced_image_ids:
                    picking_image = p
                    break
            slide_idx += 1

        # 内容页
        if texts['bullets']:
            add_content_slide(pres, texts['title'], texts['bullets'],
                              image=picking_image,
                              kicker=kicker, interactive=interactive,
                              idx=slide_idx)
            if picking_image is not None:
                produced_image_ids.add(str(picking_image))
            slide_idx += 1

        # 题目卡片页
        if texts['quizzes']:
            add_quiz_cards_slide(pres, texts['title'], texts['quizzes'], kicker=kicker)
            slide_idx += 1

        prev_kicker = kicker

    add_end_slide(pres, course_title)

    out_path = output or (course_dir / f"{course_dir.name}.pptx")
    pres.save(str(out_path))
    return out_path


def main():
    parser = argparse.ArgumentParser(
        description='TeachAny · HTML → PPTX 导出工具（v5.34.12 · 图文美观版）'
    )
    parser.add_argument('course_dir', help='课件目录，形如 examples/math-m-vertex-form')
    parser.add_argument('-o', '--output', help='输出 pptx 路径，默认 <课件目录>/<目录名>.pptx')
    args = parser.parse_args()

    course_dir = Path(args.course_dir)
    if not course_dir.exists() or not course_dir.is_dir():
        print(f"❌ 课件目录不存在: {course_dir}")
        sys.exit(1)

    try:
        out_path = export_pptx(course_dir, Path(args.output) if args.output else None)
        size_kb = out_path.stat().st_size / 1024
        print(f"✅ PPTX 导出完成: {out_path} ({size_kb:.1f} KB)")
    except Exception as e:
        import traceback
        traceback.print_exc()
        print(f"❌ 导出失败: {e}")
        sys.exit(3)


if __name__ == '__main__':
    main()
